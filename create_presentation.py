"""
FaceMate Presentation Generator
Creates a professional PowerPoint presentation for the FaceMate project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_facemate_presentation():
    """Create the FaceMate PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    PRIMARY_COLOR = RGBColor(41, 128, 185)  # Blue
    SECONDARY_COLOR = RGBColor(231, 76, 60)  # Red
    ACCENT_COLOR = RGBColor(46, 204, 113)  # Green
    TEXT_COLOR = RGBColor(44, 62, 80)  # Dark Blue
    
    def add_title_slide():
        """Slide 1: Title Slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "FaceMate"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(66)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "AI-Powered Face Recognition Attendance System"
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = TEXT_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Tagline
        tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.8))
        tagline_frame = tagline_box.text_frame
        tagline_frame.text = "Revolutionizing Attendance Management with Advanced AI"
        tagline_para = tagline_frame.paragraphs[0]
        tagline_para.font.size = Pt(20)
        tagline_para.font.italic = True
        tagline_para.font.color.rgb = SECONDARY_COLOR
        tagline_para.alignment = PP_ALIGN.CENTER
        
        # Footer
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = "Presented at Global AI Summit | February 2026"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = TEXT_COLOR
        footer_para.alignment = PP_ALIGN.CENTER
    
    def add_problem_statement():
        """Slide 2: Problem Statement"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        
        # Title
        title = slide.shapes.title
        title.text = "Problem Statement"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🎯 The Challenge: Traditional Attendance Systems Are Broken"
        
        # Add bullet points
        bullets = [
            ("Time-Consuming", "Manual roll calls waste 10-15 minutes daily"),
            ("Proxy Attendance", "Students can mark attendance for absent friends"),
            ("Inaccurate Records", "Human errors in manual entry (15-20% error rate)"),
            ("No Real-Time Insights", "Teachers can't identify at-risk students early"),
            ("Security Concerns", "Easy to manipulate attendance records")
        ]
        
        for problem, detail in bullets:
            p = tf.add_paragraph()
            p.text = f"❌ {problem}: {detail}"
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(10)
        
        # Add impact box
        impact_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
        impact_frame = impact_box.text_frame
        impact_frame.text = "📊 Annual Loss: 50+ hours per teacher wasted on attendance"
        impact_para = impact_frame.paragraphs[0]
        impact_para.font.size = Pt(20)
        impact_para.font.bold = True
        impact_para.font.color.rgb = SECONDARY_COLOR
        impact_para.alignment = PP_ALIGN.CENTER
    
    def add_solution():
        """Slide 3: Our Solution"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = "Our Solution - FaceMate"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "💡 Innovative AI-Powered System"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        
        solutions = [
            "✅ Instant Recognition: Mark attendance in < 1 second per student",
            "✅ 99.9% Accuracy: Advanced anti-spoofing prevents fraud",
            "✅ Real-Time Analytics: AI-powered insights and risk assessment",
            "✅ Zero Contact: Completely touchless and hygienic",
            "✅ 98% Storage Reduction: Revolutionary embedding-only technology"
        ]
        
        for solution in solutions:
            p = tf.add_paragraph()
            p.text = solution
            p.level = 0
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.space_before = Pt(12)
        
        # Value proposition
        value_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(7), Inches(0.8))
        value_frame = value_box.text_frame
        value_frame.text = '"From hours of manual work to seconds of AI magic"'
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(22)
        value_para.font.bold = True
        value_para.font.italic = True
        value_para.font.color.rgb = ACCENT_COLOR
        value_para.alignment = PP_ALIGN.CENTER
    
    def add_innovations():
        """Slide 4: Key Innovations"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = "Key Innovations"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "🚀 What Makes FaceMate Special"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        
        innovations = [
            ("1. Advanced Anti-Spoofing Security 🛡️", 
             "Eye blink + texture analysis | 99.9% prevention accuracy"),
            ("2. Embedding-Only Storage ⚡", 
             "6MB → 12KB per student | 99.8% storage reduction"),
            ("3. AI-Powered Analytics Engine 🧠", 
             "Risk assessment + trend prediction + auto recommendations"),
            ("4. Multi-Frame Verification 🎯", 
             "Quality assessment + adaptive thresholding | 98.5% accuracy")
        ]
        
        for title_text, detail in innovations:
            p = tf.add_paragraph()
            p.text = title_text
            p.level = 0
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
            p.space_before = Pt(10)
            
            p2 = tf.add_paragraph()
            p2.text = detail
            p2.level = 1
            p2.font.size = Pt(16)
            p2.font.color.rgb = TEXT_COLOR
    
    def add_tech_stack():
        """Slide 5: Technology Stack"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = "Technology Stack"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Left column - Backend & AI
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
        left_frame = left_box.text_frame
        
        left_content = [
            ("💻 Backend & Framework", True),
            ("🐍 Django 5.2.5 - Web framework", False),
            ("🗄️ SQLite - Database", False),
            ("🔐 Custom Auth System", False),
            ("", True),
            ("🤖 AI & Computer Vision", True),
            ("ArcFace - Face recognition", False),
            ("RetinaFace - Face detection", False),
            ("NumPy - Embeddings", False),
            ("OpenCV - Video processing", False),
            ("dlib - Facial landmarks", False)
        ]
        
        for text, is_header in left_content:
            if text:
                p = left_frame.add_paragraph() if left_frame.text else left_frame.paragraphs[0]
                if not left_frame.text:
                    left_frame.text = text
                else:
                    p.text = text
                p.font.size = Pt(16) if is_header else Pt(14)
                p.font.bold = is_header
                p.font.color.rgb = PRIMARY_COLOR if is_header else TEXT_COLOR
                p.space_before = Pt(10) if is_header else Pt(3)
        
        # Right column - Frontend & Deployment
        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
        right_frame = right_box.text_frame
        
        right_content = [
            ("🌐 Frontend", True),
            ("HTML5/CSS3 - Modern UI", False),
            ("JavaScript - Interactions", False),
            ("Bootstrap - Responsive", False),
            ("", True),
            ("☁️ Deployment", True),
            ("ngrok - Public access", False),
            ("Git - Version control", False),
            ("", True),
            ("📊 Performance", True),
            ("✅ 98.5% Accuracy", False),
            ("⚡ Real-time (30 FPS)", False),
            ("🛡️ 99.9% Anti-spoofing", False)
        ]
        
        for text, is_header in right_content:
            if text:
                p = right_frame.add_paragraph() if right_frame.text else right_frame.paragraphs[0]
                if not right_frame.text:
                    right_frame.text = text
                else:
                    p.text = text
                p.font.size = Pt(16) if is_header else Pt(14)
                p.font.bold = is_header
                p.font.color.rgb = PRIMARY_COLOR if is_header else TEXT_COLOR
                p.space_before = Pt(10) if is_header else Pt(3)
    
    def add_features():
        """Slide 6: Core Features"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = "Core Features"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Three columns
        col_width = 3
        col_height = 5
        
        # Column 1 - Developers
        dev_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(col_width), Inches(col_height))
        dev_frame = dev_box.text_frame
        dev_frame.text = "👨‍💻 For Developers"
        dev_frame.paragraphs[0].font.size = Pt(18)
        dev_frame.paragraphs[0].font.bold = True
        dev_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        dev_features = [
            "• Bulk registration",
            "• Management dashboard",
            "• Analytics & reporting",
            "• Database optimization"
        ]
        for feature in dev_features:
            p = dev_frame.add_paragraph()
            p.text = feature
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_COLOR
        
        # Column 2 - Teachers
        teacher_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.5), Inches(col_width), Inches(col_height))
        teacher_frame = teacher_box.text_frame
        teacher_frame.text = "👨‍🏫 For Teachers"
        teacher_frame.paragraphs[0].font.size = Pt(18)
        teacher_frame.paragraphs[0].font.bold = True
        teacher_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        teacher_features = [
            "• Live face attendance",
            "• Attendance history",
            "• Student management",
            "• Analytics dashboard",
            "• Manual override"
        ]
        for feature in teacher_features:
            p = teacher_frame.add_paragraph()
            p.text = feature
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_COLOR
        
        # Column 3 - Students
        student_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(col_width), Inches(col_height))
        student_frame = student_box.text_frame
        student_frame.text = "👨‍🎓 For Students"
        student_frame.paragraphs[0].font.size = Pt(18)
        student_frame.paragraphs[0].font.bold = True
        student_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        student_features = [
            "• Personal dashboard",
            "• Attendance %",
            "• History tracking",
            "• Profile management",
            "• Monthly reports"
        ]
        for feature in student_features:
            p = student_frame.add_paragraph()
            p.text = feature
            p.font.size = Pt(14)
            p.font.color.rb = TEXT_COLOR
    
    def add_team():
        """Slide 7: Team & Development"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = "Team & Development"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Content
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = "👥 Project Team Distribution"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        
        teams = [
            ("🏗️ Architecture & Backend", "Django, Database, Auth, APIs"),
            ("🤖 AI & Computer Vision", "Face recognition, Anti-spoofing, Embeddings"),
            ("🎨 Frontend & UI/UX", "Responsive design, Dashboards, Camera integration"),
            ("📊 Testing & Deployment", "Testing, Optimization, Security, Deploy")
        ]
        
        for team_name, description in teams:
            p = tf.add_paragraph()
            p.text = team_name
            p.level = 0
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = PRIMARY_COLOR
            p.space_before = Pt(8)
            
            p2 = tf.add_paragraph()
            p2.text = description
            p2.level = 1
            p2.font.size = Pt(15)
            p2.font.color.rgb = TEXT_COLOR
        
        # Timeline
        timeline_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.8))
        timeline_frame = timeline_box.text_frame
        timeline_frame.text = "⏱️ Development: 8 weeks | Core (2w) → AI (3w) → Optimize (2w) → Deploy (1w)"
        timeline_para = timeline_frame.paragraphs[0]
        timeline_para.font.size = Pt(16)
        timeline_para.font.bold = True
        timeline_para.font.color.rgb = ACCENT_COLOR
        timeline_para.alignment = PP_ALIGN.CENTER
    
    def add_impact():
        """Slide 8: Impact & Results"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title
        title = slide.shapes.title
        title.text = "Impact & Results"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        # Metrics box
        metrics_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.5))
        metrics_frame = metrics_box.text_frame
        metrics_frame.text = "📊 Performance Metrics"
        metrics_frame.paragraphs[0].font.size = Pt(20)
        metrics_frame.paragraphs[0].font.bold = True
        metrics_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        metrics = [
            "✅ Recognition: 98.5%",
            "✅ Speed: 30 FPS real-time",
            "✅ Storage: 99.8% reduction",
            "✅ Security: 99.9% anti-spoof",
            "✅ Response: < 100ms"
        ]
        
        for metric in metrics:
            p = metrics_frame.add_paragraph()
            p.text = metric
            p.font.size = Pt(16)
            p.font.color.rgb = ACCENT_COLOR
            p.font.bold = True
            p.space_before = Pt(5)
        
        # Deployment stats
        deploy_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(2.5))
        deploy_frame = deploy_box.text_frame
        deploy_frame.text = "🚀 Live Deployment"
        deploy_frame.paragraphs[0].font.size = Pt(20)
        deploy_frame.paragraphs[0].font.bold = True
        deploy_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        
        deploy_stats = [
            "📚 13 Students registered",
            "👨‍🏫 5 Teachers active",
            "🔐 9 AI Embeddings",
            "⚡ Zero security incidents",
            "⏱️ 50+ hrs saved/teacher"
        ]
        
        for stat in deploy_stats:
            p = deploy_frame.add_paragraph()
            p.text = stat
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.font.bold = True
            p.space_before = Pt(5)
        
        # Scalability table
        scale_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        scale_frame = scale_box.text_frame
        scale_frame.text = "📈 Scalability: 100 students (150MB→1.2MB) | 1,000 students (1.5GB→12MB) | 10,000 students (15GB→120MB)"
        scale_para = scale_frame.paragraphs[0]
        scale_para.font.size = Pt(16)
        scale_para.font.color.rgb = TEXT_COLOR
        scale_para.alignment = PP_ALIGN.CENTER
        
        # Future vision
        future_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1))
        future_frame = future_box.text_frame
        future_frame.text = "🎯 Future: Multi-campus | Mobile app | LMS integration | Global rollout"
        future_para = future_frame.paragraphs[0]
        future_para.font.size = Pt(18)
        future_para.font.bold = True
        future_para.font.color.rgb = SECONDARY_COLOR
        future_para.alignment = PP_ALIGN.CENTER
    
    # Create all slides
    print("Creating slide 1: Title...")
    add_title_slide()
    
    print("Creating slide 2: Problem Statement...")
    add_problem_statement()
    
    print("Creating slide 3: Solution...")
    add_solution()
    
    print("Creating slide 4: Innovations...")
    add_innovations()
    
    print("Creating slide 5: Tech Stack...")
    add_tech_stack()
    
    print("Creating slide 6: Features...")
    add_features()
    
    print("Creating slide 7: Team...")
    add_team()
    
    print("Creating slide 8: Impact...")
    add_impact()
    
    # Save presentation
    filename = "FaceMate_Presentation.pptx"
    prs.save(filename)
    print(f"\n✅ Presentation created successfully: {filename}")
    print(f"📊 Total slides: 8")
    print(f"📁 Location: {filename}")
    
    return filename

if __name__ == "__main__":
    create_facemate_presentation()
